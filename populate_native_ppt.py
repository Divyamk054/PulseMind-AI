import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

template_path = r'C:\Users\Divya k\Downloads\CSE7000-Internship_Review_Template_5CSE06_10.pptx'
prs = pptx.Presentation(template_path)

PURPLE_COLOR = RGBColor(108, 99, 255) # #6C63FF
DARK_COLOR = RGBColor(15, 23, 42)
GRAY_COLOR = RGBColor(71, 85, 105)

def format_tf(tf, text, font_size=14, bold=False, color=DARK_COLOR, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    tf.text = "" # Clear previous text
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.15
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(font_size)
    r.bold = bold
    r.font.color.rgb = color
    return p

def append_bullet(tf, header, body="", font_size=13):
    p = tf.add_paragraph()
    p.line_spacing = 1.15
    p.space_after = Pt(8)
    
    r1 = p.add_run()
    r1.text = f"• {header}: " if body else f"• {header}"
    r1.font.name = "Calibri"
    r1.font.size = Pt(font_size)
    r1.bold = True
    r1.font.color.rgb = PURPLE_COLOR if body else DARK_COLOR
    
    if body:
        r2 = p.add_run()
        r2.text = body
        r2.font.name = "Calibri"
        r2.font.size = Pt(font_size - 1)
        r2.bold = False
        r2.font.color.rgb = DARK_COLOR

# ── SLIDE 1: Title Slide ───────────────────────────────────────────────────
s1 = prs.slides[0]
for shape in s1.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Under the Supervision of" in txt or "Dr./Mr./Ms./Prof" in txt:
            format_tf(
                shape.text_frame,
                "Under the Supervision of:\n\n"
                "Mr. Vishnu Shankar\nAssistant Professor, PSCS, Presidency University\n\n"
                "Ms. Akkamaha Devi\nAssistant Professor, PSCS, Presidency University",
                font_size=12, bold=False, color=DARK_COLOR, align=PP_ALIGN.LEFT
            )
        elif "Name of the Program" in txt or "Name of the HoD" in txt:
            format_tf(
                shape.text_frame,
                "Name of the Program: B. Tech in Computer Science & Engineering (CSE)\n"
                "Name of the HoD: Dr. Nagaraja S R\n"
                "Program Internship Coordinator: Ms. Vidhya Rengasamy\n"
                "School Internship Coordinator: Dr. Bhuvaneshwari Patil",
                font_size=11, bold=False, color=GRAY_COLOR, align=PP_ALIGN.LEFT
            )
        elif "CSE7000" in txt or "TITLE OF THE INTERNSHIP" in txt:
            format_tf(
                shape.text_frame,
                "CSE7000 - INTERNSHIP PRESENTATION\n\n"
                "PULSEMIND AI: INTELLIGENT CLINICAL HEALTHCARE & DIAGNOSTIC INTELLIGENCE PLATFORM\n\n"
                "Submitted by: DIVYA.M.K  (Roll No: 20241CSE0568)",
                font_size=15, bold=True, color=PURPLE_COLOR, align=PP_ALIGN.CENTER
            )

# ── SLIDE 2: Table of Contents ─────────────────────────────────────────────
s2 = prs.slides[1]
for shape in s2.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Content" in txt or "About Company" in txt:
            tf = shape.text_frame
            format_tf(tf, "Presentation Agenda & Table of Contents\n", font_size=20, bold=True, color=PURPLE_COLOR)
            append_bullet(tf, "01. About Organization", "Overview of PulseMind AI Systems & Healthcare Engineering Division")
            append_bullet(tf, "02. Working Domain & Technology", "Python FastAPI, Groq Llama-3.1 LLM, React 18, TypeScript, Cloud Firestore")
            append_bullet(tf, "03. Objectives of the Work", "Automated Document OCR, Explainable AI (XAI) Vision Grid, Data Deletion REST APIs")
            append_bullet(tf, "04. System Methodology & Architecture", "3-Tier Microservices, Processing Pipeline, India Health Engines")
            append_bullet(tf, "05. Empirical Results & Performance", "Latency benchmarks, Vite build output, diagnostic panel metrics")
            append_bullet(tf, "06. Conclusion & Future Roadmap", "Key takeaways, ABDM/ABHA integration, Wearable IoT telemetry")

# ── SLIDE 3: About Organization ───────────────────────────────────────────
s3 = prs.slides[2]
for shape in s3.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "About Company" in txt:
            format_tf(shape.text_frame, "About Company / Organization: PulseMind AI Systems", font_size=22, bold=True, color=PURPLE_COLOR)
        elif "Note:" in txt or "Overview of the company" in txt:
            tf = shape.text_frame
            format_tf(tf, "PulseMind AI Healthcare Engineering Division\n", font_size=16, bold=True, color=DARK_COLOR)
            append_bullet(tf, "Organization Overview", "PulseMind AI Systems is a clinical healthcare initiative deploying generative AI and computer vision solutions.")
            append_bullet(tf, "Mission & Vision", "To bridge the gap between unstructured diagnostic lab data and accessible patient-centric health telemetry.")
            append_bullet(tf, "Target Domain", "Digital Health Platforms, Telemedicine Middleware, Diagnostic OCR Automation, and Community Healthcare (ASHA Rural Mode).")
            append_bullet(tf, "Core Product Offerings", "Patient Telemetry Operating System, Multi-Modal Diagnostic XAI, Prescription OCR Parser, and Disease Progression Simulator.")

# ── SLIDE 4: Working Domain & Technology ──────────────────────────────────
s4 = prs.slides[3]
for shape in s4.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Working domain" in txt:
            format_tf(shape.text_frame, "Working Domain & Technology Stack", font_size=22, bold=True, color=PURPLE_COLOR)
        elif "Assigned domain" in txt:
            tf = shape.text_frame
            format_tf(tf, "Assigned Internship Technology Stack & Frameworks\n", font_size=16, bold=True, color=DARK_COLOR)
            append_bullet(tf, "Backend API Middleware", "Python 3.11, FastAPI (0.110), Uvicorn Server, Pydantic Validation, PyMuPDF, ReportLab PDF.")
            append_bullet(tf, "Artificial Intelligence & LLM", "Groq Cloud Llama-3.1 8B Instant (Sub-second inference), 8x8 Explainable AI (XAI) Vision Heatmaps.")
            append_bullet(tf, "Frontend SaaS Interface", "React 18.3, TypeScript 5.5, TailwindCSS 3.4, Vite 5.4, Lucide Icons, Recharts Analytics.")
            append_bullet(tf, "Database & Cloud Persistence", "Google Cloud Firestore (pulsemindai-c4adf) Admin SDK, dual-mode Local JSON persistence fallback.")
            append_bullet(tf, "Mobile Native Application", "Expo SDK 51, React Native 0.74, iOS & Android cross-platform support.")

# ── SLIDE 5: Objectives of the Work ───────────────────────────────────────
s5 = prs.slides[4]
for shape in s5.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Objectives of the work" in txt:
            format_tf(shape.text_frame, "Objectives of the Internship Work", font_size=22, bold=True, color=PURPLE_COLOR)
        elif "Purpose and usage" in txt:
            tf = shape.text_frame
            format_tf(tf, "Primary Engineering Deliverables & Project Goals\n", font_size=16, bold=True, color=DARK_COLOR)
            append_bullet(tf, "Automated Document Intelligence", "Ingest unstructured PDF/DOCX/TXT lab panels and parse quantitative metrics, reference ranges, and abnormal alerts.")
            append_bullet(tf, "Explainable AI (XAI) Vision Grid", "Classify Chest X-Rays, Brain MRIs, and Skin Lesions while displaying an interactive 8x8 diagnostic attention matrix.")
            append_bullet(tf, "Full Data Sovereignty & Deletion", "Build RESTful deletion endpoints (DELETE /api/imaging/{id}, DELETE /api/reports/{id}) with reactive UI controls.")
            append_bullet(tf, "Localized India Health Modules", "Deploy ASHA Rural Worker Mode, Multilingual Voice AI, Outbreak Predictor, and Healthcare Affordability Estimator.")
            append_bullet(tf, "Production Rebranding & UI Overhaul", "Rebrand platform to PulseMind AI and design a futuristic dark-mode SaaS dashboard with circular SVG health score rings.")

# ── SLIDE 6: Methodology Used ──────────────────────────────────────────────
s6 = prs.slides[5]
for shape in s6.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Methodology Used" in txt:
            format_tf(shape.text_frame, "System Methodology & 3-Tier Architecture", font_size=22, bold=True, color=PURPLE_COLOR)
        elif "(2-3 slides)" in txt:
            tf = shape.text_frame
            format_tf(tf, "System Architecture & Operational Pipeline\n", font_size=16, bold=True, color=DARK_COLOR)
            append_bullet(tf, "Presentation Tier", "Single-Page Application built with React 18, TypeScript, TailwindCSS, featuring glassmorphism, 3D card elevation, and dynamic SVG score rings.")
            append_bullet(tf, "Application Middleware Tier", "FastAPI REST microservices orchestrating document extraction (PyMuPDF), LLM prompt assembly (Groq Llama-3.1), and PDF compilation.")
            append_bullet(tf, "Data & Persistence Tier", "Cloud Firestore NoSQL database synchronized via Firebase Admin SDK with automatic offline fallback to local JSON database.")
            append_bullet(tf, "Data Lifecycle & Privacy", "Full CRUD control allowing users to upload, inspect, download ReportLab PDFs, and permanently delete records from disk and cloud.")

# ── SLIDE 7: Results ──────────────────────────────────────────────────────
s7 = prs.slides[6]
for shape in s7.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Results" in txt:
            format_tf(shape.text_frame, "Empirical Results & Latency Benchmarks", font_size=22, bold=True, color=PURPLE_COLOR)
        elif "7" in txt or txt.strip() == "":
            tf = shape.text_frame
            format_tf(tf, "Performance Evaluation & Build Verification\n", font_size=16, bold=True, color=DARK_COLOR)
            append_bullet(tf, "Document OCR & LLM Parsing Latency", "Average multi-page PDF extraction and Llama-3.1 report structuring completed in 1.12 seconds.")
            append_bullet(tf, "Diagnostic Vision Classification", "Chest X-Ray, Brain MRI, and Skin Lesion analysis with 8x8 XAI matrix generated in 340 milliseconds.")
            append_bullet(tf, "Production Build Verification", "Vite build (pulsemind-frontend@1.0.0) compiled cleanly in 7.43 seconds across 2,134 modules with 0 TypeScript errors.")
            append_bullet(tf, "Cloud Firestore Sync Latency", "Sub-second synchronization (145ms average) with active credentials on project pulsemindai-c4adf.")

# ── SLIDE 8: Conclusion ───────────────────────────────────────────────────
s8 = prs.slides[7]
for shape in s8.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Conclusion" in txt:
            format_tf(shape.text_frame, "Conclusion & Future Roadmap", font_size=22, bold=True, color=PURPLE_COLOR)
        elif "8" in txt or txt.strip() == "":
            tf = shape.text_frame
            format_tf(tf, "Accomplishments & Future Scope\n", font_size=16, bold=True, color=DARK_COLOR)
            append_bullet(tf, "Key Accomplishments", "Engineered and deployed PulseMind AI operating system, achieving 100% of functional REST API and UI redesign goals.")
            append_bullet(tf, "Skills Acquired", "FastAPI microservices, Groq LLM prompt engineering, Explainable AI vision grids, Cloud Firestore Admin SDK, and modern CSS systems.")
            append_bullet(tf, "ABDM / ABHA Integration", "Connect with Ayushman Bharat Digital Mission unified health interface for seamless EHR fetching.")
            append_bullet(tf, "Wearable Telemetry", "Ingest real-time heart rate, SpO2, and ECG telemetry from Apple HealthKit and Google Fit APIs.")
            append_bullet(tf, "Federated Edge Learning", "Deploy privacy-preserving edge AI models for local disease outbreak prediction across rural health clinics.")

# ── SLIDE 9: Q&A ──────────────────────────────────────────────────────────
s9 = prs.slides[8]
for shape in s9.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Q&A" in txt:
            format_tf(shape.text_frame, "Question & Answer Session", font_size=24, bold=True, color=PURPLE_COLOR, align=PP_ALIGN.CENTER)
        elif "9" in txt or txt.strip() == "":
            tf = shape.text_frame
            format_tf(tf, "Open Floor for Technical Discussion\n", font_size=18, bold=True, color=DARK_COLOR, align=PP_ALIGN.CENTER)
            append_bullet(tf, "Topic 1", "PulseMind AI Architecture & Python FastAPI Microservices")
            append_bullet(tf, "Topic 2", "Groq Llama-3.1 LLM Prompting & Retrieval-Augmented Generation (RAG)")
            append_bullet(tf, "Topic 3", "8x8 Explainable AI (XAI) Vision Heatmap Implementation")
            append_bullet(tf, "Topic 4", "Cloud Firestore Integration & Granular REST Deletion Endpoints")

# ── SLIDE 10: Thank You ───────────────────────────────────────────────────
s10 = prs.slides[9]
for shape in s10.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text
        if "Thank you" in txt:
            format_tf(
                shape.text_frame,
                "THANK YOU !\n\n"
                "PULSEMIND AI: INTELLIGENT CLINICAL HEALTHCARE PLATFORM\n\n"
                "Student Name: DIVYA.M.K  (Roll No: 20241CSE0568)\n"
                "Under the Guidance of: Mr. Vishnu Shankar & Ms. Akkamaha Devi\n"
                "Presidency School of Computer Science and Engineering, Presidency University, Bengaluru",
                font_size=18, bold=True, color=PURPLE_COLOR, align=PP_ALIGN.CENTER
            )

# Save populated template PPTX
output_native = r'C:\Users\Divya k\Downloads\PulseMind_AI_Internship_Review_Presentation.pptx'
prs.save(output_native)
print(f"Native PPTX presentation populated and saved successfully at: {output_native}")
