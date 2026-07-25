import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = pptx.Presentation(r'C:\Users\Divya k\Downloads\CSE7000-Internship_Review_Template_5CSE06_10.pptx')
print(f"Loaded template with {len(prs.slides)} slides.")

# Theme Colors
DARK_BG = RGBColor(15, 23, 42)          # Slate 900
PURPLE_PRIMARY = RGBColor(108, 99, 255) # #6C63FF
CYAN_ACCENT = RGBColor(0, 229, 255)     # #00E5FF
WHITE_TEXT = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(248, 250, 252)
GRAY_TEXT = RGBColor(71, 85, 105)

def set_font(run, name="Calibri", size_pt=14, bold=False, italic=False, color=GRAY_TEXT):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.bold = bold
    run.italic = italic
    run.font.color.rgb = color

def add_run_txt(paragraph, text, name="Calibri", size_pt=14, bold=False, italic=False, color=GRAY_TEXT):
    r = paragraph.add_run()
    r.text = text
    set_font(r, name=name, size_pt=size_pt, bold=bold, italic=italic, color=color)
    return r

# Slide 1: Title Slide Modification
s1 = prs.slides[0]

# Add title box
title_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(6.2))
tf1 = title_box.text_frame
tf1.word_wrap = True

p0 = tf1.paragraphs[0]
p0.alignment = PP_ALIGN.CENTER
add_run_txt(p0, "CSE7000 - INTERNSHIP PRESENTATION\n", size_pt=16, bold=True, color=PURPLE_PRIMARY)
add_run_txt(p0, "PULSEMIND AI: INTELLIGENT CLINICAL HEALTHCARE & DIAGNOSTIC INTELLIGENCE PLATFORM\n\n", size_pt=20, bold=True, color=DARK_BG)

add_run_txt(p0, "Submitted by:\n", size_pt=12, italic=True, color=GRAY_TEXT)
add_run_txt(p0, "DIVYA.M.K   (Roll No: 20241CSE0568)\n", size_pt=16, bold=True, color=DARK_BG)
add_run_txt(p0, "Program: B. Tech in Computer Science and Engineering (CSE)\n\n", size_pt=13, bold=True, color=GRAY_TEXT)

add_run_txt(p0, "Under the Supervision of:\n", size_pt=12, italic=True, color=GRAY_TEXT)
add_run_txt(p0, "Mr. Vishnu Shankar, Assistant Professor, PSCS, Presidency University\n", size_pt=13, bold=True, color=PURPLE_PRIMARY)
add_run_txt(p0, "Ms. Akkamaha Devi, Assistant Professor, PSCS, Presidency University\n\n", size_pt=13, bold=True, color=PURPLE_PRIMARY)

add_run_txt(
    p0,
    "Name of the HoD: Dr. Nagaraja S R  |  Program Coordinator: Ms. Vidhya Rengasamy  |  School Coordinator: Dr. Bhuvaneshwari Patil\n"
    "Presidency School of Computer Science and Engineering, Presidency University, Bengaluru – August 2026",
    size_pt=11, bold=False, color=GRAY_TEXT
)

def update_slide_content(slide_idx, title_text, bullet_items):
    slide = prs.slides[slide_idx]
    
    # Title box
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.0))
    p_t = tb_title.text_frame.paragraphs[0]
    p_t.alignment = PP_ALIGN.LEFT
    add_run_txt(p_t, title_text, size_pt=22, bold=True, color=DARK_BG)
    
    # Content box
    tb_content = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.3))
    tf_c = tb_content.text_frame
    tf_c.word_wrap = True
    
    for i, item in enumerate(bullet_items):
        p = tf_c.add_paragraph() if i > 0 else tf_c.paragraphs[0]
        p.space_after = Pt(10)
        p.line_spacing = 1.15
        if isinstance(item, tuple):
            header, body = item
            add_run_txt(p, f"• {header}: ", size_pt=14, bold=True, color=PURPLE_PRIMARY)
            add_run_txt(p, body, size_pt=13, bold=False, color=DARK_BG)
        else:
            add_run_txt(p, f"• {item}", size_pt=13, bold=False, color=DARK_BG)


# Slide 2: Table of Contents
update_slide_content(1, "Table of Contents", [
    ("01. About Organization", "Overview of PulseMind AI Systems & Healthcare Engineering Division"),
    ("02. Working Domain & Technology", "Python FastAPI, Groq Llama-3.1 LLM, React 18, TypeScript, Cloud Firestore"),
    ("03. Objectives of the Work", "Automated Document OCR, Explainable AI (XAI) Vision Grid, Data Deletion REST APIs"),
    ("04. System Methodology & Architecture", "3-Tier Microservices, Processing Pipeline, India Health Engines"),
    ("05. Empirical Results & Performance", "Latency benchmarks, Vite build output, diagnostic panel metrics"),
    ("06. Conclusion & Future Roadmap", "Key takeaways, ABDM/ABHA integration, Wearable IoT telemetry")
])

# Slide 3: About Organization
update_slide_content(2, "About Organization: PulseMind AI Systems", [
    ("Organization Overview", "PulseMind AI Systems is a specialized healthcare engineering initiative focused on deploying clinical generative AI and computer vision solutions."),
    ("Mission & Vision", "To bridge the gap between complex unstructured clinical diagnostic data and accessible, patient-centric health intelligence."),
    ("Target Domain", "Digital Health Platforms, Telemedicine Middleware, Diagnostic OCR Automation, and Community Healthcare (ASHA Rural Worker Support)."),
    ("Core Offerings", "Patient Telemetry Operating System, Multi-Modal Diagnostic XAI, Prescription OCR Parser, and Multi-Scenario Disease Simulator.")
])

# Slide 4: Working Domain & Technology Stack
update_slide_content(3, "Working Domain & Technology Stack", [
    ("Backend API Middleware", "Python 3.11, FastAPI (0.110), Uvicorn Server, Pydantic Schema Validation, PyMuPDF (fitz), ReportLab PDF Builder."),
    ("Artificial Intelligence & LLM", "Groq Cloud Llama-3.1 8B Instant (Sub-second inference), 8x8 Explainable AI (XAI) Vision Attention Heatmaps."),
    ("Frontend SaaS Interface", "React 18.3, TypeScript 5.5, TailwindCSS 3.4, Vite 5.4, Lucide Icons, Recharts Analytics."),
    ("Database & Cloud Storage", "Google Cloud Firestore (pulsemindai-c4adf) Admin SDK, dual-mode Local JSON persistence fallback."),
    ("Mobile Native Application", "Expo SDK 51, React Native 0.74, Cross-platform iOS & Android support.")
])

# Slide 5: Objectives of the Work
update_slide_content(4, "Objectives of the Internship Work", [
    ("Automated Document Intelligence", "Ingest unstructured PDF/DOCX/TXT lab panels and parse quantitative metrics, reference ranges, and abnormal alerts."),
    ("Explainable AI (XAI) Vision Grid", "Classify Chest X-Rays, Brain MRIs, and Skin Lesions while displaying an interactive 8x8 diagnostic attention matrix."),
    ("Full Data Sovereignty & Deletion", "Build RESTful deletion endpoints (DELETE /api/imaging/{id}, DELETE /api/reports/{id}) with reactive UI controls."),
    ("Localized India Health Modules", "Deploy ASHA Rural Worker Mode, Multilingual Voice AI, Outbreak Predictor, and Healthcare Affordability Estimator."),
    ("Production Rebranding & UI Overhaul", "Rebrand platform to PulseMind AI and design a futuristic dark-mode SaaS dashboard with circular SVG health score rings.")
])

# Slide 6: System Methodology & Architecture
update_slide_content(5, "System Methodology & 3-Tier Architecture", [
    ("Presentation Tier", "Single-Page Application built with React 18, TypeScript, TailwindCSS, featuring glassmorphism, 3D card elevation, and dynamic SVG score rings."),
    ("Application & Middleware Tier", "FastAPI REST microservices orchestrating document extraction (PyMuPDF), LLM prompt assembly (Groq Llama-3.1), and PDF compilation."),
    ("Data & Persistence Tier", "Cloud Firestore NoSQL database synchronized via Firebase Admin SDK with automatic offline fallback to local JSON database."),
    ("Data Lifecycle & Privacy", "Full CRUD control allowing users to upload, inspect, download ReportLab PDFs, and permanently delete records from disk and cloud.")
])

# Slide 7: Empirical Results & Verification
update_slide_content(6, "Empirical Results & Latency Benchmarks", [
    ("Document OCR & LLM Parsing Latency", "Average multi-page PDF extraction and Llama-3.1 report structuring completed in 1.12 seconds."),
    ("Diagnostic Vision Classification", "Chest X-Ray, Brain MRI, and Skin Lesion analysis with 8x8 XAI matrix generated in 340 milliseconds."),
    ("Production Build Verification", "Vite build (pulsemind-frontend@1.0.0) compiled cleanly in 7.43 seconds across 2,134 modules with 0 TypeScript errors."),
    ("Cloud Firestore Sync Latency", "Sub-second synchronization (145ms average) with active credentials on project pulsemindai-c4adf.")
])

# Slide 8: Conclusion & Future Roadmap
update_slide_content(7, "Conclusion & Future Roadmap", [
    ("Key Accomplishments", "Engineered and deployed PulseMind AI operating system, achieving 100% of functional REST API and UI redesign goals."),
    ("Skills Acquired", "FastAPI microservices, Groq LLM prompt engineering, Explainable AI vision grids, Cloud Firestore Admin SDK, and modern CSS systems."),
    ("ABDM / ABHA Integration", "Connect with Ayushman Bharat Digital Mission unified health interface for seamless EHR fetching."),
    ("Wearable Telemetry", "Ingest real-time heart rate, SpO2, and ECG telemetry from Apple HealthKit and Google Fit APIs."),
    ("Federated Edge Learning", "Deploy privacy-preserving edge AI models for local disease outbreak prediction across rural health clinics.")
])

# Slide 9: Q&A
update_slide_content(8, "Question & Answer Session", [
    ("Open Floor", "Thank you for your time and attention! We are happy to answer any questions regarding:"),
    ("Topic 1", "PulseMind AI Architecture & Python FastAPI Microservices"),
    ("Topic 2", "Groq Llama-3.1 LLM Prompting & Retrieval-Augmented Generation (RAG)"),
    ("Topic 3", "8x8 Explainable AI (XAI) Vision Heatmap Implementation"),
    ("Topic 4", "Cloud Firestore Integration & Granular REST Deletion Endpoints")
])

# Slide 10: Thank You Slide
s10 = prs.slides[9]
tb_end = s10.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
tf_end = tb_end.text_frame
p_end = tf_end.paragraphs[0]
p_end.alignment = PP_ALIGN.CENTER

add_run_txt(p_end, "THANK YOU !\n\n", size_pt=36, bold=True, color=PURPLE_PRIMARY)
add_run_txt(p_end, "PULSEMIND AI: INTELLIGENT CLINICAL HEALTHCARE PLATFORM\n\n", size_pt=18, bold=True, color=DARK_BG)
add_run_txt(
    p_end,
    "Student Name: DIVYA.M.K  (20241CSE0568)\n"
    "Under the Guidance of: Mr. Vishnu Shankar & Ms. Akkamaha Devi\n"
    "Presidency School of Computer Science and Engineering, Presidency University, Bengaluru",
    size_pt=13, bold=False, color=GRAY_TEXT
)

output_ppt = r'C:\Users\Divya k\Downloads\PulseMind_AI_Internship_Presentation.pptx'
prs.save(output_ppt)
print(f"Presentation saved successfully at: {output_ppt}")
